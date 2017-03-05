## プレセンターションクラス

import pptx
import abc
import gpt.NoteText

# プレセンテーションクラスのインタフェース
class IPresentation:
	__metaclass__ = abc.ABCMeta

	@abc.abstractmethod
	def getNoteTexts(self):
		raise NotImplementedError()

	@abc.abstractmethod
	def getPlatformPresentation(self):
		raise NotImplementedError()


class Presentation(IPresentation):
	def __init__(self, filename):
		self.LoadedPresentation = pptx.Presentation(filename)
		self.NoteTexts = list()
		for p, slide in enumerate(self.LoadedPresentation.slides):
			if slide.has_notes_slide:
				note = slide.notes_slide
				page_number = p + 1
				note_text = str()
				for paragraph in note.notes_text_frame.paragraphs:
					note_text += paragraph.text + '\n'
				self.NoteTexts.append(gpt.NoteText.NoteText(page_number, note_text))

	def getNoteTexts(self):
		return self.NoteTexts

	def getPlatformPresentation(self):
		return self.LoadedPresentation
